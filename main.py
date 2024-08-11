import google.generativeai as genai
import os
import sys
from docx_parser import RAGFlowDocxParser
from excel_parser import RAGFlowExcelParser
from pdf_parser import RAGFlowPdfParser
from pptx import Presentation

# Configure the Gemini API
genai.configure(api_key=os.environ["API_KEY"])


def summarize_text(text):
    # Initialize the model
    model = genai.GenerativeModel('gemini-1.5-flash')

    # Build a prompt that explicitly requests a summary of no more than 200 words and identifies the main tag and secondary tag that represent the core topics of the document.
    prompt = (
        "Please summarize the following document in no more than 200 words. "
        "Additionally, identify one main tag and one secondary tag that represent the core topics of the document.\n\n"
        f"{text}\n\n"
        "Summary (max 200 words):\n"
        "Main Tag:\n"
        "Secondary Tag:"
    )
    response = model.generate_content(prompt)
    return response.text


def main(file_path):
    # Select parser based on file type
    if file_path.endswith(".docx"):
        parser = RAGFlowDocxParser()
        sections, tables = parser(file_path)
        # Join the extracted text into a single string
        all_text = " ".join([section[0] for section in sections])
    elif file_path.endswith(".xlsx"):
        parser = RAGFlowExcelParser()
        all_text = " ".join(parser(file_path))
    elif file_path.endswith(".pdf"):
        parser = RAGFlowPdfParser()
        sections, tables = parser(file_path)
        all_text = " ".join(sections)
    elif file_path.endswith(".pptx"):
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        all_text = " ".join(text)
    else:
        raise ValueError(
            "Unsupported file type. Please provide a .docx, .xlsx, .pdf, or .pptx file."
        )

    # Print the first and last 1000 characters of the text to verify extraction
    # print("First 1000 characters:")
    # print(all_text[:1000])
    # print("\n\nLast 1000 characters:")
    # print(all_text[-1000:])

    # Call the Gemini API to get the summary
    summary = summarize_text(all_text)

    print("Summary:")
    print(summary)


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python main.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    main(file_path)
